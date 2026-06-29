"""
PPT 생성 엔진
파싱된 슬라이드 데이터를 회사 표준 템플릿 기반 .pptx로 변환
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from config import (
    Font, Color, Layout,
    CoverPos, SectionPos, ContentPos, TocPos,
    ROMAN_NUMERALS,
)


def generate_pptx(
    slides_data: list[dict],
    template_path: str | Path,
    output_path: str | Path,
):
    """
    슬라이드 데이터 리스트를 받아 PPT 파일을 생성한다.

    Args:
        slides_data: parser.parse_docx()의 반환값
        template_path: 회사 표준 템플릿 .pptx 경로
        output_path: 출력 .pptx 경로
    """
    prs = Presentation(str(template_path))

    # 섹션 번호 카운터 (로마자 자동 부여)
    section_counter = 0
    current_section_title = ""

    # 섹션 제목 목록 수집 (목차 자동 생성용)
    section_titles = []
    for s in slides_data:
        if s.get("type") == "섹션헤더":
            section_titles.append(s.get("제목", ""))

    # 슬라이드 생성
    for slide_data in slides_data:
        slide_type = slide_data.get("type", "본문")

        if slide_type == "커버":
            _create_cover(prs, slide_data)
            # 커버 직후 목차 자동 삽입
            if section_titles:
                _create_toc(prs, section_titles)

        elif slide_type == "섹션헤더":
            section_counter += 1
            numeral = ROMAN_NUMERALS[section_counter - 1] if section_counter <= len(ROMAN_NUMERALS) else str(section_counter)
            current_section_title = f"{numeral}. {slide_data.get('제목', '')}"
            _create_section_header(prs, slide_data, numeral)

        elif slide_type == "본문":
            _create_content(prs, slide_data, current_section_title)

        elif slide_type == "목차":
            # 수동 목차 (자동 생성과 별개로 강사가 직접 지정한 경우)
            custom_items = slide_data.get("내용", section_titles)
            _create_toc(prs, custom_items)

    prs.save(str(output_path))


# === 슬라이드 생성 함수 ===

def _create_cover(prs: Presentation, data: dict):
    """커버 슬라이드 생성"""
    slide = _add_slide(prs, Layout.COVER)

    # 메인 타이틀
    title = data.get("제목", "강의 제목")
    _add_textbox(
        slide,
        left=CoverPos.TITLE_LEFT,
        top=CoverPos.TITLE_TOP,
        width=CoverPos.TITLE_WIDTH,
        height=CoverPos.TITLE_HEIGHT,
        text=title,
        font_name=Font.TITLE_NAME,
        font_size=Font.TITLE_SIZE,
        alignment=PP_ALIGN.CENTER,
        font_color=Color.TITLE_RGB,
    )

    # 부제목 (있는 경우)
    subtitle = data.get("부제목", "")
    if subtitle:
        _add_textbox(
            slide,
            left=CoverPos.SUBTITLE_LEFT,
            top=CoverPos.SUBTITLE_TOP,
            width=CoverPos.SUBTITLE_WIDTH,
            height=CoverPos.SUBTITLE_HEIGHT,
            text=subtitle,
            font_name=Font.BODY_NAME,
            font_size=Pt(18),
            alignment=PP_ALIGN.CENTER,
            font_color=Color.BODY_RGB,
        )


def _create_section_header(prs: Presentation, data: dict, numeral: str):
    """섹션헤더 슬라이드 생성"""
    slide = _add_slide(prs, Layout.COVER)  # 커버와 동일 레이아웃

    title = f"{numeral}. {data.get('제목', '')}"
    _add_textbox(
        slide,
        left=SectionPos.TITLE_LEFT,
        top=SectionPos.TITLE_TOP,
        width=SectionPos.TITLE_WIDTH,
        height=SectionPos.TITLE_HEIGHT,
        text=title,
        font_name=Font.TITLE_NAME,
        font_size=Font.TITLE_SIZE,
        alignment=PP_ALIGN.CENTER,
        font_color=Color.TITLE_RGB,
    )


def _create_content(prs: Presentation, data: dict, section_title: str):
    """본문 슬라이드 생성"""
    slide = _add_slide(prs, Layout.BLANK)

    # 1. 섹션 표시 바 (상단 좌측)
    if section_title:
        _add_textbox(
            slide,
            left=ContentPos.SECTION_BAR_LEFT,
            top=ContentPos.SECTION_BAR_TOP,
            width=ContentPos.SECTION_BAR_WIDTH,
            height=ContentPos.SECTION_BAR_HEIGHT,
            text=section_title,
            font_name=Font.SECTION_BAR_NAME,
            font_size=Font.SECTION_BAR_SIZE,
            font_color=Color.SECTION_BAR_RGB,
            vertical_anchor=MSO_ANCHOR.MIDDLE,
        )

    # 2. 소제목 (섹션 바 아래)
    subtitle = data.get("소제목", "")
    if subtitle:
        _add_textbox(
            slide,
            left=ContentPos.SUBTITLE_LEFT,
            top=ContentPos.SUBTITLE_TOP,
            width=ContentPos.SUBTITLE_WIDTH,
            height=ContentPos.SUBTITLE_HEIGHT,
            text=subtitle,
            font_name=Font.SUBTITLE_NAME,
            font_size=Font.SUBTITLE_SIZE,
            font_color=Color.TITLE_RGB,
        )

    # 3. 본문 콘텐츠 (불릿 포인트)
    bullets = data.get("내용", [])
    if bullets:
        _add_bullet_textbox(
            slide,
            left=ContentPos.BODY_LEFT,
            top=ContentPos.BODY_TOP,
            width=ContentPos.BODY_WIDTH,
            height=ContentPos.BODY_HEIGHT,
            bullets=bullets,
        )


def _create_toc(prs: Presentation, section_titles: list[str]):
    """목차 슬라이드 생성 (섹션 제목 리스트 기반)"""
    slide = _add_slide(prs, Layout.BLANK)

    # "< 목 차 >" 타이틀
    _add_textbox(
        slide,
        left=TocPos.TITLE_LEFT,
        top=TocPos.TITLE_TOP,
        width=TocPos.TITLE_WIDTH,
        height=TocPos.TITLE_HEIGHT,
        text="< 목 차 >",
        font_name=Font.BODY_NAME,
        font_size=Font.TOC_TITLE_SIZE,
        alignment=PP_ALIGN.CENTER,
        font_color=Color.BODY_RGB,
    )

    # 섹션 제목 리스트
    if section_titles:
        txBox = slide.shapes.add_textbox(
            TocPos.LIST_LEFT,
            TocPos.LIST_TOP,
            TocPos.LIST_WIDTH,
            TocPos.LIST_HEIGHT,
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, title in enumerate(section_titles):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()

            numeral = ROMAN_NUMERALS[i] if i < len(ROMAN_NUMERALS) else str(i + 1)
            para.text = f"{numeral}. {title}"
            para.space_after = Pt(14)

            run = para.runs[0]
            run.font.name = Font.BODY_NAME
            run.font.size = Font.TOC_ITEM_SIZE
            run.font.color.rgb = RGBColor(*Color.BODY_RGB)


# === 유틸리티 함수 ===

def _add_slide(prs: Presentation, layout_index: int):
    """지정 레이아웃으로 빈 슬라이드 추가"""
    layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(layout)


def _add_textbox(
    slide,
    left, top, width, height,
    text: str,
    font_name: str = Font.BODY_NAME,
    font_size=Font.BODY_SIZE,
    font_color: tuple = Color.BODY_RGB,
    bold: bool = False,
    alignment=None,
    vertical_anchor=None,
):
    """텍스트박스를 추가하고 스타일을 적용한다."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    if vertical_anchor is not None:
        tf.paragraphs[0].alignment = alignment
        txBox.text_frame.auto_size = None

    para = tf.paragraphs[0]
    para.text = text

    if alignment is not None:
        para.alignment = alignment

    run = para.runs[0]
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*font_color)

    if vertical_anchor is not None:
        tf.paragraphs[0].alignment = alignment


def _add_bullet_textbox(
    slide,
    left, top, width, height,
    bullets: list[str],
):
    """불릿 포인트 텍스트박스를 추가한다. 불릿 수에 따라 폰트 크기 자동 조절."""
    # 불릿 수에 따른 폰트 크기 결정
    if len(bullets) <= 5:
        font_size = Font.BODY_SIZE        # 18pt
    elif len(bullets) <= 8:
        font_size = Font.BODY_SIZE_SMALL  # 16pt
    else:
        font_size = Font.BODY_SIZE_TINY   # 14pt

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()

        para.text = f"• {bullet}"
        para.space_after = Pt(8)

        run = para.runs[0]
        run.font.name = Font.BODY_NAME
        run.font.size = font_size
        run.font.color.rgb = RGBColor(*Color.BODY_RGB)
